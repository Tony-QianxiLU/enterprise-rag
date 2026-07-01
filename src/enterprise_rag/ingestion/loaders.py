"""Document loading: turns raw uploaded bytes into a LoadedDocument.

Dispatch is by filename extension. Each format's parser errors are caught and
re-raised as ValueError so callers (the API layer) can treat all parse
failures uniformly as a 4xx-worthy user error rather than a 500.
"""

from io import BytesIO

from docx import Document as DocxDocument
from pptx import Presentation
from pypdf import PdfReader
from pypdf.errors import PdfReadError

from enterprise_rag.schemas import DocumentType, LoadedDocument

_EXTENSION_TO_TYPE: dict[str, DocumentType] = {
    ".pdf": DocumentType.PDF,
    ".docx": DocumentType.DOCX,
    ".txt": DocumentType.TXT,
    ".md": DocumentType.MARKDOWN,
    ".markdown": DocumentType.MARKDOWN,
    ".pptx": DocumentType.PPTX,
}


def detect_document_type(filename: str) -> DocumentType:
    lower_name = filename.lower()
    for extension, document_type in _EXTENSION_TO_TYPE.items():
        if lower_name.endswith(extension):
            return document_type

    supported = ", ".join(sorted(_EXTENSION_TO_TYPE))
    raise ValueError(f"Unsupported file type for '{filename}'. Supported types: {supported}")


def _load_pdf(file_bytes: bytes) -> str:
    try:
        reader = PdfReader(BytesIO(file_bytes))
        pages = [page.extract_text() or "" for page in reader.pages]
    except (PdfReadError, ValueError) as error:
        raise ValueError("Could not read PDF content") from error

    return "\n\n".join(page.strip() for page in pages if page.strip())


def _load_docx(file_bytes: bytes) -> str:
    try:
        document = DocxDocument(BytesIO(file_bytes))
        paragraphs = [paragraph.text for paragraph in document.paragraphs]
    except Exception as error:
        raise ValueError("Could not read DOCX content") from error

    return "\n".join(paragraph.strip() for paragraph in paragraphs if paragraph.strip())


def _load_text(file_bytes: bytes) -> str:
    try:
        return file_bytes.decode("utf-8").strip()
    except UnicodeDecodeError as error:
        raise ValueError("Could not decode text content as UTF-8") from error


def _load_pptx(file_bytes: bytes) -> str:
    try:
        presentation = Presentation(BytesIO(file_bytes))
        texts: list[str] = []
        for slide in presentation.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = shape.text_frame.text.strip()
                    if text:
                        texts.append(text)
    except Exception as error:
        raise ValueError("Could not read PPTX content") from error

    return "\n".join(texts)


def load_document(file_bytes: bytes, filename: str) -> LoadedDocument:
    document_type = detect_document_type(filename)

    if document_type is DocumentType.PDF:
        text = _load_pdf(file_bytes)
    elif document_type is DocumentType.DOCX:
        text = _load_docx(file_bytes)
    elif document_type in (DocumentType.TXT, DocumentType.MARKDOWN):
        text = _load_text(file_bytes)
    else:
        text = _load_pptx(file_bytes)

    return LoadedDocument(filename=filename, document_type=document_type, text=text)
