from io import BytesIO
from pathlib import Path

import pytest
from docx import Document as DocxDocument
from pptx import Presentation
from pypdf import PdfWriter

from enterprise_rag.ingestion.loaders import detect_document_type, load_document
from enterprise_rag.schemas import DocumentType


def test_detect_document_type_for_supported_extensions() -> None:
    assert detect_document_type("report.pdf") == DocumentType.PDF
    assert detect_document_type("report.docx") == DocumentType.DOCX
    assert detect_document_type("notes.txt") == DocumentType.TXT
    assert detect_document_type("notes.md") == DocumentType.MARKDOWN
    assert detect_document_type("notes.markdown") == DocumentType.MARKDOWN
    assert detect_document_type("slides.pptx") == DocumentType.PPTX


def test_detect_document_type_rejects_unsupported_extension() -> None:
    with pytest.raises(ValueError, match="Unsupported file type"):
        detect_document_type("archive.zip")


def test_load_document_txt_decodes_utf8_content() -> None:
    document = load_document("Hello RAG".encode("utf-8"), "notes.txt")

    assert document.filename == "notes.txt"
    assert document.document_type == DocumentType.TXT
    assert document.text == "Hello RAG"


def test_load_document_markdown_decodes_utf8_content() -> None:
    document = load_document("# Heading\n\nSome text".encode("utf-8"), "notes.md")

    assert document.document_type == DocumentType.MARKDOWN
    assert document.text == "# Heading\n\nSome text"


def test_load_document_rejects_unsupported_extension() -> None:
    with pytest.raises(ValueError, match="Unsupported file type"):
        load_document(b"content", "archive.zip")


def test_load_document_docx_round_trips_paragraph_text(tmp_path: Path) -> None:
    docx_path = tmp_path / "fixture.docx"
    docx_document = DocxDocument()
    docx_document.add_paragraph("Hello DOCX world")
    docx_document.save(docx_path)

    document = load_document(docx_path.read_bytes(), "fixture.docx")

    assert document.document_type == DocumentType.DOCX
    assert "Hello DOCX world" in document.text


def test_load_document_pptx_round_trips_text_frame_text(tmp_path: Path) -> None:
    pptx_path = tmp_path / "fixture.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    text_box = slide.shapes.add_textbox(0, 0, 1000000, 1000000)
    text_box.text_frame.text = "Hello PPTX world"
    presentation.save(pptx_path)

    document = load_document(pptx_path.read_bytes(), "fixture.pptx")

    assert document.document_type == DocumentType.PPTX
    assert "Hello PPTX world" in document.text


def test_load_document_pdf_loads_blank_page_without_raising() -> None:
    writer = PdfWriter()
    writer.add_blank_page(width=200, height=200)
    buffer = BytesIO()
    writer.write(buffer)

    document = load_document(buffer.getvalue(), "fixture.pdf")

    assert document.document_type == DocumentType.PDF
    assert document.filename == "fixture.pdf"
    assert isinstance(document.text, str)
